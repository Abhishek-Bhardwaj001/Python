# Library imports
import pandas as pd
import fitz  # PyMuPDF
import pdfplumber
import base64
from io import BytesIO
from PIL import Image
import io

# DOCX related
from docx import Document as Docx_loader
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.ns import qn

# Excel
from openpyxl import load_workbook

# PowerPoint
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Standard libraries
import zipfile
import os
import shutil
import uuid
from pathlib import Path
import subprocess
import sys
# OLE files (older Office formats)
import olefile
from Python.Gen_AI.rag_systems.utils.data_transform_helpers import install_libreoffice
class TextExtractor:
    def __init__(self):
        pass

    #=====================*:Main Function Call for Text Extractor:*=========================    
    def convert_to_markdown(self, data, page_metadata=False,verbose=False):
        """
        Extract content from various text document types (Microsoft Documents, Text files, .msg files) and convert it into markdown format. The function automatically detects the type of the document provided.

        Args:
        data (str | bytes | bytearray): Document file to extract text from. Can be provided as a binary object or a file path.
        page_metadata (bool, optional): Flag to include content metadata in the output. Metadata may include data types (Text, Image, Table), page number, document title, etc. Defaults to False.
        verbose (bool, optional): Flag to enable consolse prints of operations performed during extraction. Defaults to False.

        Returns:
        str | dict: Extracted markdown content in string format if `page_metadata` is False, or in dictionary format if `page_metadata` is True.
        """
        if verbose:
            print("---Entered markdown---")        
        text = [] if page_metadata else ""

        ext = self._detect_file_type(data,verbose)
        if verbose:
            print(f"------The Detected Extension of the file is: {ext}------")

        if ext in ['doc','ppt','xls','docx']:  # Converting Legacy Microsoft Documents to thier new formats
            if ext.lower()=="doc":
                output_format="docx"        
            elif ext.lower()=="ppt":
                output_format="pptx"
            elif ext.lower()=="xls":
                output_format="xlsx"
            elif ext.lower()=="docx":
                output_format="pdf"    
            data = self._document_conversion(data,input_format=ext,output_format=output_format,verbose=verbose)
            ext = self._detect_file_type(data,verbose)
            print(f"The New extension of the file is:{ext}")

        if ext.lower() == "pdf":
            text = self._extract_pdf(data,page_metadata,verbose)

        elif ext.lower() in ["docx","docm"]:
            text = self._extract_docx(data,page_metadata,verbose)
        
        elif ext.lower() in ["pptx","pptm"]:
            text = self._extract_pptx(data,page_metadata,verbose)
        
        elif ext.lower() in ["xlsx","xlsm"]:
            text = self._extract_xlsx(data,page_metadata,verbose)
        else:
            print(f"The Class Vectorization does not support {ext} file type")

        if verbose:
            print("------Exiting Markdown------")
        if text:
            return text
        else:
            raise ValueError("------No text found in the document------")

    #=====================*:Main Function Call for Text Extractor Ends here:*========================= 

    #=====================*:Private Methods of the Text Extractor Class:*========================= 
    def _detect_file_type(self, data,verbose=False):
        """
        Detects the type of a document file based on the provided data. This function can handle both file paths (as strings) and binary data, and it attempts to identify the file type (such as PDF, Word, Excel, PowerPoint, etc.) using various detection methods.

        Args:
            data (str | bytes | bytearray): The document file to inspect, either as a path (string) or binary data.
            verbose (bool, optional): If set to True, enables prints debug information to the console. Defaults to False.

        Returns:
            str: A string representing the detected file type. Possible return values include:
                - 'pdf': Portable Document Format
                - 'doc': Microsoft Word (legacy format)
                - 'docx': Microsoft Word (XML-based format)
                - 'docm': Microsoft Word Macro-Enabled (XML-based format)
                - 'dotm': Microsoft Word Macro-Enabled Template
                - 'xls': Microsoft Excel (legacy format)
                - 'xlsx': Microsoft Excel (XML-based format)
                - 'xlsm': Microsoft Excel Macro-Enabled (XML-based format)
                - 'ppt': Microsoft PowerPoint (legacy format)
                - 'pptx': Microsoft PowerPoint (XML-based format)
                - 'pptm': Microsoft PowerPoint Macro-Enabled (XML-based format)
                - 'txt': Plain text file
                - 'msg': Outlook message format
                - 'Unknown': If the file type cannot be determined.

        Raises:
            Exception: Multiple exception handling is implemented for file reading and type detection processes, which can raise specific errors based on the detection mechanism used.

        Example Function Call:
            file_type = self._detect_file_type("document.docx", verbose=True)
        """
        if verbose:
                print("---Entered detect_file_type---")
            
        if isinstance(data, str):
            try:
                return Path(data).suffix[1:].lower()
            except Exception as e:
                if verbose:
                    print(f"Invalid file path: {e}")

        elif isinstance(data, (bytes, bytearray)):
            if verbose:
                print("------The Detected data is binary------")
            
            try:
                # Check for binary file signatures
                if data.startswith(b'%PDF-'):
                    return "pdf"
                
                elif data.startswith(b'PK\x03\x04'):
                    try:
                        with zipfile.ZipFile(io.BytesIO(data)) as zip_file:
                            namelist = zip_file.namelist()
                            # Detect Word files  
                            if any(name.startswith("word/") for name in namelist):
                                if any(".bin" in name for name in namelist):
                                    return "docm"
                                elif any(name.endswith(".xml") for name in namelist):
                                    return "docx"
                                elif any(name.endswith(".dotm") for name in namelist):
                                    return "dotm"
                                else:
                                    print("Failed to detect Word file type.")

                            # Detect PowerPoint files
                            elif any(name.startswith("ppt/") for name in namelist):
                                if any(".bin" in name for name in namelist):
                                    return "pptm"
                                else:
                                    return "pptx"

                            # Detect Excel files
                            elif any(name.startswith("xl/") for name in namelist):
                                if any(".bin" in name for name in namelist):
                                    return "xlsm"
                                else:
                                    return "xlsx"
                    except Exception as e:
                        if verbose:
                            print(f"\033[32mFailed to detect the file type of binary zip data: {e}!\033[0m")

                elif data.startswith(b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1\x00\x00'):
                    if verbose:
                        print("--------Enter older file version Detection block -----")
                    try:
                        legacy_doc = BytesIO(data)
                        with olefile.OleFileIO(legacy_doc) as ole:
                            if verbose:
                                print("entered ole Detection")
                            if ole.exists("WordDocument"):
                                return "doc"
                            elif ole.exists("PowerPoint Document"):
                                return "ppt"
                            elif ole.exists("Workbook") or ole.exists("Book"):  # Check for Excel files
                                return "xls"
                            elif ole.exists("__nameid_version1.0"):
                                return "msg"
                            else:
                                return "Unknown"
                        if verbose:
                            print("-------------------Exiting OLE File detection-----------------------")
                    except Exception as e:
                        if verbose:
                            print(f"Failed to detect the file type of OLE file: {e}")
                    
                else:
                    try:
                        decode = data.decode("utf-8")  # Ignore decode errors
                        sample = decode[:200]
                        return "txt"
                    except Exception as e:
                        if verbose:
                            print(f"Failed to detect the file type of binary data: {e}")

                    raise ValueError('Failed to detect the file type of data')
            except Exception as e:
                if verbose:
                    print(f"Detection failed for binary data with error: {e}")
                raise ValueError('Failed to detect the file type of data')
        if verbose:
            print("---------Exiting File Type Detection---------")
    

    def _extract_pdf(self,data,page_metadata=False,verbose=False):
        """
        Extracts text and tables from a PDF document and converts them into markdown format. 
        The function can handle input as either binary data or a file path.

        Args:
            data (str | bytes | bytearray): The PDF document to extract text and tables from, specified as a file path or binary data.
            page_metadata (bool, optional): If True, includes metadata for each extracted content item (text/table), such as the document title, data type, page number, and dimensions. Defaults to False.
            verbose (bool, optional): If set to True, enables cpnsole print of operations and prints debug information to the console. Defaults to False.
        
        Returns:
            str | dictionary | None: A list of Dictionaries containing extracted content in markdown format. Each item may include text or tables, along with their respective metadata if 
                        `page_metadata` is True. Returns None if extraction fails.

        Example:
            extracted_content = self._extract_pdf("document.pdf", page_metadata=True, verbose=True)

        Raises:
            Exception: Raises an exception if the extraction process encounters an error, for debugging purposes.
        """
        if verbose:
            print("---Extracting Text from PDF---")

        if isinstance(data, (bytes, bytearray)):
            pdf_stream = BytesIO(data)
            fitz_doc = fitz.open(stream=pdf_stream, filetype="pdf")
            plumber_doc = pdfplumber.open(pdf_stream)
        else:
            fitz_doc = fitz.open(data)
            plumber_doc = pdfplumber.open(data)

        doc_metadata = fitz_doc.metadata
        text_extract = [] if page_metadata else ""
        
        try:
            if fitz_doc:
                for page in range(len(fitz_doc)):
                    fitz_page = fitz_doc.load_page(page)
                    plumber_page = plumber_doc.pages[page]
                    page_width = round(fitz_page.rect.width, 2)
                    page_height = round(fitz_page.rect.height, 2)
                    page_number = page + 1

                    # Extract tables and their bounding boxes
                    tables = plumber_page.extract_tables()
                    table_bboxes = [table.bbox for table in plumber_page.find_tables()]

                    # Extract text blocks
                    page_text = ""
                    for block in fitz_page.get_text("blocks"):
                        block_rect = fitz.Rect(block[:4])
                        is_in_table = False
                        for bbox in table_bboxes:
                            table_rect = fitz.Rect(bbox)
                            if table_rect.intersects(block_rect):
                                is_in_table = True
                                break
                        if not is_in_table:
                            page_text += block[4].strip() + " "

                    page_text = page_text.strip()

                    # Add cleaned text
                    if page_metadata:
                        if page_text:
                            text_extract.append({
                                "content": page_text,
                                "metadata": {
                                    "pdf_title": doc_metadata.get("title", ""),
                                    "Data_type": "text",
                                    "page_Number": page_number,
                                    "page_width": page_width,
                                    "page_height": page_height
                                }
                            })
                    else:
                        text_extract += page_text + "\n"

                    # Add tables
                    for table_index, table in enumerate(tables):
                        header = [str(cell) if cell is not None else "" for cell in table[0]]
                        table_markdown = "| " + "|".join(header) + "|\n"
                        table_markdown += "| " + "|".join(["---"] * len(header)) + "|\n"
                        for row in table[1:]:
                            table_markdown += "| " + "|".join([str(cell) if cell is not None else "" for cell in row]) + "|\n"

                        if page_metadata:
                            text_extract.append({
                                "content": table_markdown,
                                "metadata": {
                                    "pdf_title": doc_metadata.get("title", ""),
                                    "Data_type": "table",
                                    "page_Number": page_number,
                                    "page_width": page_width,
                                    "page_height": page_height,
                                    "table_index": table_index + 1,
                                    "rows": len(table) - 1,
                                    "columns": len(header)
                                }
                            })
                        else:
                            text_extract += table_markdown + "\n"

                    # == Image data conerversion Code ==
                    images = fitz_page.get_images(full=True)
                    for img_index, image in enumerate(images):
                        image_id = image[0]
                        base_image = fitz_doc.extract_image(image_id)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]

                        image_pil = Image.open(BytesIO(image_bytes))
                        buffered = BytesIO()
                        image_pil.save(buffered, format=image_ext.upper())
                        base64_encrypted = base64.b64encode(buffered.getvalue()).decode("utf-8")
                        if page_metadata:
                            text_extract.append({
                                "content": base64_encrypted,
                                "metadata":{
                                    "pdf_title": doc_metadata.get("title",""),
                                    "Data_type":"image",
                                    "page_Number": page_number,
                                    "page_width": page_width,
                                    "page_height": page_height,
                                    'image_index': img_index+1,
                                    'image_ext': image_ext,
                                    'image_size': f'{image_pil.size}'
                                }
                            })
                        else:
                            text_extract += ("").join(base64_encrypted)

            fitz_doc.close()
            plumber_doc.close()
            if verbose:
                print("---Text Extraction Complete---")
            return text_extract
        except Exception as e:
            print(f"failed to extract text from pdf due to {e}")
            return None

    def _extract_docx(self, data, page_metadata=False, verbose=False):
        """
        Extracts text and tables from a DOCX document and converts them into markdown format. The function can handle input as either binary data or a file path.

        Args:
            data (str | bytes | bytearray): The DOCX document to extract text and tables from, specified as a file path or binary data.
            page_metadata (bool, optional): If True, includes metadata for each extracted content item (text/table), such as data type and block index. Defaults to False.
            verbose (bool, optional): If set to True, enables console prints of operations and prints debug information to the console. Defaults to False.

        Returns:
            list | str: A list containing extracted content in markdown format if `page_metadata` is True. If `page_metadata` is False, returns a string with extracted content concatenated together. Returns None if extraction fails.

        Raises:
            ValueError: Raises a ValueError if the input type is neither a string, bytes, nor 
                        BytesIO, along with a error message.

        Example:
            extracted_content = self._extract_docx("document.docx", page_metadata=True, verbose=True)
        """
        if verbose:
            print("--------Entered Doc to text Conversion------------")

        if isinstance(data, (bytes, bytearray)):
            doc = Docx_loader(BytesIO(data))
        elif isinstance(data, str):
            doc = Docx_loader(data)

        else:
            raise ValueError("Unsupported input type. Please provide a str, bytes, or BytesIO.")

        if page_metadata:
            chunks = []
        else:
            chunks = ""
        block_index = 0
        parent_elm = doc.element.body
        
        try:
            if doc:
                for child in parent_elm.iterchildren():
                    if isinstance(child, CT_P):
                        para = Paragraph(child, doc)
                        try:
                            if page_metadata:
                                if para.text.strip():
                                    chunks.append({
                                        "content": para.text.strip(),
                                        "metadata": {
                                            "Data_type": "text",
                                            "block_index": block_index
                                        }
                                    })
                                    block_index += 1
                            else:
                                chunks += ("").join(para.text.strip())
                                chunks += ("\n\n")
                        except Exception as e:
                            print(f"[WARN] Failed to extract text: {e}")


                        for run in para.runs:
                            if run._element.xpath('.//*[local-name()="drawing"]'):
                                for blip in run._element.xpath('.//*[local-name()="blip"]'):
                                    rId = blip.get(qn('r:embed'))
                                    if not rId:
                                        continue
                                    image_part = doc.part.related_parts[rId]
                                    image_bytes = image_part.blob
                                    try:
                                        # Handle only supported image types, skip WMF
                                        if image_part.content_type == 'image/wmf':
                                            print("[WARN] WMF images are not supported and will be skipped.")
                                            continue  # Skip WMF images

                                        # Proceed with other supported image formats
                                        image_pil = Image.open(BytesIO(image_bytes))
                                        buffered = BytesIO()
                                        image_pil.save(buffered, format=image_pil.format or "PNG")
                                        image_base64 = base64.b64encode(buffered.getvalue()).decode()
                                        
                                        if page_metadata:
                                            chunks.append({
                                                "content": image_base64,
                                                "metadata": {
                                                    "Data_type": "image",
                                                    "block_index": block_index,
                                                    "width": image_pil.width,
                                                    "height": image_pil.height,
                                                    "image_format": image_pil.format or "PNG"
                                                }
                                            })
                                            block_index += 1
                                        else:
                                            chunks += ("").join(image_base64)
                                            chunks += ("\n")
                                    except Exception as e:
                                        print(f"[WARN] Failed to extract image: {e}")
                                        continue  # Continue on failure for this image

                    elif isinstance(child, CT_Tbl):
                        table = Table(child, doc)
                        rows = table.rows
                        if rows:
                            headers = [cell.text.strip() for cell in rows[0].cells]
                            markdown = "| " + " | ".join(headers) + " |\n"
                            markdown += "| " + " | ".join(["---"] * len(headers)) + " |\n"

                            for row in rows[1:]:
                                values = [cell.text.strip() for cell in row.cells]
                                markdown += "| " + " | ".join(values) + " |\n"
                            if page_metadata:
                                chunks.append({
                                    "content": markdown.strip(),
                                    "metadata": {
                                        "Data_type": "table",
                                        "block_index": block_index,
                                        "rows": len(rows) - 1,
                                        "columns": len(headers)
                                    }
                                })
                                block_index += 1
                            else:
                                chunks += ("").join(markdown.strip())
                                chunks += ("\n\n")
                
                if verbose:
                    print("--------Exiting Doc to text Conversion------------")                
                return chunks
        except Exception as e:
            print(f"Failed to extract text from doc due to: {e}")
            return None
      
    

    def _extract_xlsx(self,data,page_metadata=False,verbose=False):
        """
        Extracts text and tables from an XLSX document and converts them into markdown format. The function can handle inputs as either binary data or a file path.

        Args:
            data (str | bytes | bytearray): The XLSX document to extract text and tables from, specified as a file path or binary data.
            page_metadata (bool, optional): If True, includes metadata for each extracted content item (text/table), such as data type, source, sheet name, and block index. Defaults to False.
            verbose (bool, optional): If set to True, enables console print of operations and prints debug information to the console. Defaults to False.

        Returns:
            str | None: A string containing extracted content in markdown format if `page_metadata` is False. If `page_metadata` is True, returns a list of dictionaries 
                        with content and metadata. Returns None if extraction fails.

        Raises:
            Exception: Raises an exception if the extraction process encounters an error.

        Call Example:
            extracted_content = self._extract_xlsx("document.xlsx", page_metadata=True, verbose=True)
        """
        if verbose:
            print("--------Entered Xlsx to text Conversion------------")
            
        if isinstance(data, (bytes,bytearray)):
            xlsx_io = BytesIO(data)
            source = "in-memory.xlsx"
        else:
            xlsx_io = data
            source = data

        if page_metadata:
            chunks = []
        else:
            chunks = ""
        block_index = 0

        # Track used rows per sheet to avoid duplicate text
        table_row_ranges = {}

        # Step 1: Extract tables using pandas
        wb = load_workbook(filename=xlsx_io, data_only=True)
        try:
            if wb:
                sheet_names=wb.sheetnames
                for sheet_name in sheet_names:
                    df = pd.read_excel(xlsx_io, sheet_name=sheet_name, engine="openpyxl")
                    if df.empty:
                        continue

                    markdown = df.to_markdown(index=False)
                    if page_metadata:
                        chunks.append({
                            "content": markdown,
                            "metadata": {
                                "Data_type": "table",
                                "source": source,
                                "sheet": sheet_name,
                                "block_index": block_index,
                                "rows": df.shape[0],
                                "columns": df.shape[1]
                            }
                        })
                    else:
                        chunks += ("").join(markdown)
                        chunks += ("\n\n")

                    table_row_ranges[sheet_name] = set(range(1, df.shape[0] + 2))  # pandas reads from row 2 onwards
                    block_index += 1

                # Step 2: Extract images and text using openpyxl
                if isinstance(data, bytes):
                    xlsx_io.seek(0)

                wb = load_workbook(filename=xlsx_io)
                for sheet in wb.worksheets:
                    used_rows = table_row_ranges.get(sheet.title, set())

                    for row in sheet.iter_rows(min_row=1):
                        row_index = row[0].row
                        if row_index in used_rows:
                            continue  # Skip rows already parsed as table

                        row_text = " | ".join([str(cell.value).strip() for cell in row if cell.value])
                        if row_text.strip():
                            if page_metadata:
                                chunks.append({
                                    "content": row_text.strip(),
                                    "metadata": {
                                        "Data_type": "text",
                                        "source": source,
                                        "sheet": sheet.title,
                                        "block_index": block_index
                                    }
                                })
                                block_index += 1
                            else:
                                chunks += ("").join(row_text.strip())
                                chunks += ("\n\n")

                    for image in sheet._images:
                        try:
                            if hasattr(image.ref, "_blob"):
                                img_bytes = image.ref._blob
                            elif isinstance(image.ref, BytesIO):
                                img_bytes = image.ref.getvalue()
                            else:
                                raise TypeError("Unsupported image reference type")

                            image_pil = Image.open(BytesIO(img_bytes))
                            buffered = BytesIO()
                            image_pil.save(buffered, format=image_pil.format or "PNG")
                            image_base64 = base64.b64encode(buffered.getvalue()).decode()
                            if page_metadata:
                                chunks.append({
                                    "content": image_base64,
                                    "metadata": {
                                        "Data_type": "image",
                                        "source": source,
                                        "sheet": sheet.title,
                                        "block_index": block_index,
                                        "width": image_pil.width,
                                        "height": image_pil.height,
                                        "image_format": image_pil.format or "PNG"
                                    }
                                })
                                block_index += 1
                            else:
                                chunks += ("").join(image_base64)
                                chunks += ("\n\n")

                        except Exception as e:
                            print(f"[WARN] Failed to extract image from sheet '{sheet.title}': {e}")
                print("-------Started with Replacing Spaces and nan--------")
                chunks=chunks.replace(" ","")
                chunks=chunks.replace("nan"," ")
                if verbose:
                    print("--------Exiting Xlsx to text Conversion------------")
                return chunks
        except Exception as e:
            print(f"[ERROR] Failed to extract tables from xlsx file: {e}")
            return None
 
    def _extract_pptx(self,data,page_metadata=False,verbose=False):
        """
        Extracts text and tables from a PPTX presentation and converts them into markdown format. The function can handle inputs as either binary data or a file path.

        Args:
            data (str | bytes | bytearray): The PPTX document to extract text and tables from, specified as a file path or binary data.
            page_metadata (bool, optional): If True, includes metadata for each extracted content item (text/table), such as data type, source, slide number, and shape index. Defaults to False.
            verbose (bool, optional): If set to True, enables consoloe print of operations and prints debug information to the console. Defaults to False.

        Returns:
            str | None: A string containing extracted content in markdown format if `page_metadata` is False. If `page_metadata` is True, returns a list of dictionaries 
                        with content and metadata. Returns None if extraction fails.

        Raises:
            Exception: Raises an exception if the extraction process encounters an error, which is verboseged accordingly.

        Function Call Example:
            extracted_content = self._extract_pptx("presentation.pptx", page_metadata=True, verbose=True)
        """
        if verbose:
            print("--------Entered Pptx to text Conversion------------")
        if isinstance(data, (bytes,bytearray)):
            pptx_io = BytesIO(data)
            source = "in-memory.pptx"
        else:
            pptx_io = data
            source = data
        try:
            if pptx_io:
                prs = Presentation(pptx_io)
                if page_metadata:
                    chunks = []
                else:
                    chunks = ""
                block_index = 0

                for slide_num, slide in enumerate(prs.slides, start=1):
                    for shape_idx, shape in enumerate(slide.shapes):
                        shape_type = shape.shape_type

                        # 📝 Text extraction
                        if shape.has_text_frame:
                            text = shape.text.strip()
                            if text:
                                if page_metadata:
                                    chunks.append({
                                        "content": text,
                                        "metadata": {
                                            "Data_type": "text",
                                            "source": source,
                                            "slide": slide_num,
                                            "shape_index": shape_idx,
                                            "block_index": block_index
                                        }
                                    })
                                    block_index += 1
                                else:
                                    chunks += ("").join(text)
                                    chunks += ("\n\n")

                        elif shape_type == MSO_SHAPE_TYPE.TABLE:
                            try:
                                table = shape.table
                                rows = list(table.rows)

                                if not rows:
                                    continue

                                # Safely get column count from the longest row
                                col_count = max((len(row.cells) for row in rows if row.cells), default=0)
                                row_count = len(rows)

                                if col_count == 0:
                                    continue  # corrupt table

                                # Try to extract header row
                                header_row = rows[0].cells[:col_count]
                                headers = [cell.text.strip() if cell.text.strip() else f"Column {i+1}" for i, cell in enumerate(header_row)]

                                markdown = "| " + " | ".join(headers) + " |\n"
                                markdown += "| " + " | ".join(["---"] * len(headers)) + " |\n"

                                # Body rows
                                for row in rows[1:]:
                                    if not row.cells:
                                        continue
                                    values = [cell.text.strip() for cell in row.cells[:col_count]]
                                    # pad missing cells
                                    if len(values) < col_count:
                                        values += [""] * (col_count - len(values))
                                    markdown += "| " + " | ".join(values) + " |\n"
                                if page_metadata:
                                    chunks.append({
                                        "content": markdown.strip(),
                                        "metadata": {
                                            "Data_type": "table",
                                            "source": source,
                                            "slide": slide_num,
                                            "shape_index": shape_idx,
                                            "block_index": block_index,
                                            "rows": row_count,
                                            "columns": col_count
                                        }
                                    })
                                    block_index += 1
                                else:
                                    chunks += ("").join(markdown.strip())
                                    chunks += ("\n\n")

                            except Exception as e:
                                print(f"[WARN] Failed to extract table on slide {slide_num}: {e}")



                        # 🖼️ Image extraction
                        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image_blob = shape.image.blob
                                image_format = shape.image.ext.upper()
                                normalized_format = "JPEG" if image_format == "JPG" else image_format

                                image_pil = Image.open(BytesIO(image_blob))
                                buffered = BytesIO()
                                image_pil.save(buffered, format=normalized_format)
                                image_base64 = base64.b64encode(buffered.getvalue()).decode()
                                if page_metadata:
                                    chunks.append({
                                        "content": image_base64,
                                        "metadata": {
                                            "Data_type": "image",
                                            "source": source,
                                            "slide": slide_num,
                                            "shape_index": shape_idx,
                                            "block_index": block_index,
                                            "width": image_pil.width,
                                            "height": image_pil.height,
                                            "image_format": normalized_format
                                        }
                                    })
                                    block_index += 1
                                else:
                                    chunks += ("").join(image_base64)
                                    chunks += ("\n\n")
                            except Exception as e:
                                print(f"[WARN] Failed to extract image on slide {slide_num}: {e}")
                if verbose:
                    print("--------Exiting Pptx to text Conversion------------")
                return chunks
            
        except Exception as e:
            print(f"[ERROR] Failed to extract text from slide {slide_num}: {e}")
            return None
        
    
    def _document_conversion(self,doc,input_format,output_format,verbose=False, libreoffice_path="/usr/bin/libreoffice"):
        """
        Converts document files from one format to another using LibreOffice in headless mode.This function writes the binary input to a temporary file, performs the conversion, and returns the converted document in binary format.

        Args:
            doc (bytes): The binary content of the document to be converted.
            input_format (str): The format of the input document (e.g., 'doc', 'docx', 'odt').
            output_format (str): The desired format for the output document (e.g., 'pdf', 'docx', 'odt').
            verbose (bool, optional): If set to True, enables prints debug information 
                                to the console. Defaults to False.
            libreoffice_path (str, optional): The file path to the LibreOffice executable. 
                                            Defaults to "/usr/bin/libreoffice".

        Returns:
            bytes: The binary content of the converted document.

        Raises:
            FileNotFoundError: If the conversion fails and the output file is not created.
            subprocess.CalledProcessError: If the LibreOffice subprocess fails.

        Example:
            converted_doc = self._document_conversion(doc_binary, 'docx', 'pdf', verbose=True)

        Notes:
            This function automatically checks for the presence of LibreOffice and will attempt to install it if not found. Temporary files are created during conversion and removed afterward.
        """
        try:
            subprocess.run(["libreoffice","--version"],capture_output=True,text=True)
            if verbose:
                print("------Libreoffice Found-----")
        except Exception as e:
            if verbose:
                print(f"------{e}-----")
                print('------Libreoffice Not Found Performing Installation------')
            install_libreoffice(verbose=verbose)

        # Create temp working directory
        temp_dir = f"/tmp/{input_format}_extract_{uuid.uuid4().hex}"
        os.makedirs(temp_dir, exist_ok=True)
        if verbose==True:
            print("----------Converting legacy Document to acceptable format-------")
        try:
            # Step 1: Write binary to .doc file
            if isinstance(doc, (bytes,bytearray)):
                doc_path = os.path.join(temp_dir, f"input.{input_format}")
                with open(doc_path, "wb") as f:
                    f.write(doc)
            elif isinstance(doc, str):
                with open(doc,'rb') as file:
                    doc_binary = file.read()
                doc_path = os.path.join(temp_dir, f"input.{input_format}")
                with open(doc_path, "wb") as f:
                    f.write(doc_binary)
                
            if verbose:
                print(f"Debug: Writing to {doc_path}, size={len(doc)}")
            subprocess.run([
                libreoffice_path, "--headless", "--convert-to", output_format, doc_path,
                "--outdir", temp_dir
            ], check=True)

            data = os.path.join(temp_dir, f"input.{output_format}")
            if not os.path.exists(data):
                raise FileNotFoundError(f"LibreOffice conversion failed — no {output_format} file created.")

            else:
                with open(data,"rb") as f:
                    data_binary=f.read()
                    if verbose:
                        print("--------Document Conversion Successful---------")
                    return data_binary

        finally:
            # Clean up temp files if you don't want them stored
            shutil.rmtree(temp_dir, ignore_errors=True)
        if verbose:
            print("------Finished Converting legacy Document to acceptable format-------")
        
